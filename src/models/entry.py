from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pathlib import Path
from models.protein_model.human_protein import HumanProtein
from models.annotation import Annotation
from dataclasses import dataclass, field
from models.image import Img

@dataclass
class Entry:
    """
    Represents a protein passport entry (pptx).

    Attributes:
        template_path (str): Path of the template ppt.
        human (HumanProtein): HumanProtein of this Entry.
        orthologs (list): List of Orthologs of this Entry.
        user_name (str): User's name.
        powerpoint (Presentation): Presentation object of this Entry.
        slides (Slides): Slides of this Entry's Presentation.
        table_cells (list): List containing text to fill table cells of first slide in this Entry.
        output_path (Path): Output path.
    """
    template_path: str
    human: HumanProtein
    orthologs: list
    user_name: str
    powerpoint: Presentation = field(init=False)
    slides: list = field(init=False)
    table_cells: list = field(init=False)
    output_path: Path = field(init=False)

    def __post_init__(self):
        """
        Post init method for Entry. Sets powerpoint, slides, output_path, and table_cells fields.
        """
        self.powerpoint = Presentation(self.template_path)
        self.slides = self.powerpoint.slides
        self.output_path = Path(__file__).parent.parent.parent / f"output_{self.human.name}" / f"{self.human.name}_protein_passport.pptx"
        self._build_table_cells()
        self._set_footer()

    def _build_table_cells(self):
        """
        Sets table_cells field.
        """
        self.table_cells = [
            [
                "Target Name: " + self.human.passport_table_data["rec_name"],
                "Aliases: " + ", ".join(self.human.passport_table_data["aliases"]),
                f"(Gene id: {self.human.name}, UniProtKB - {self.human.id})"
            ],
            [self.human.passport_table_data["target_type"]],
            [""],
            [
                f"{self.human.passport_table_data['length']} aa {self.human.passport_table_data['mass']} kDa",
                ""
            ],
            [f"{o.organism.value[0]}: %" for o in self.orthologs],
            [
                f"Experimental PDBs: {', '.join(self.human.passport_table_data['exp_pdbs'])}",
                f"Predicted: {self.human.pred_pdb_id}"
            ],
            [f"{self.human.passport_table_data['exp_pattern']}."],
            [f"{self.human.passport_table_data['known_activity']}."]
        ]

    def _set_footer(self):
        """
        Writes user's name to footer.
        """
        for slide in self.slides:
            for shape in slide.shapes:
                if 'Footer' in shape.name:
                    shape.text=self.user_name
                    break

    def populate_info_table_slide(self, img: Img):
        """
        Populates the first slide of protein passport ppt template.

        Args:
            img (Img): Human protein 3d structure image.
        """
        slide = self.slides[0]

        for shape in slide.shapes:
            if 'Title' in shape.name:
                title = shape
            if 'Picture' in shape.name:
                picture = shape
            if 'Text' in shape.name:
                pbd_id_caption = shape
            if shape.has_table:
                table = shape.table
        
        if title:
            title.text = "Protein Passport - " + self.human.name

        if picture:
            img.vertical()
            picture.insert_picture(img.path)
        
        if pbd_id_caption:
            pbd_id_caption.text = img.caption
            pbd_id_caption.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        
        if table:
            for i, cell_text in enumerate(self.table_cells):
                cell = table.cell(i, 1)
                cell.text = "\n".join(cell_text)

                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rbg = RGBColor(0, 0, 0)

        self.powerpoint.save(self.output_path)
    
    def populate_hu_seq_slide(self):
        """
        Populates the second slide of protein passport ppt template.
        """
        slide = self.slides[1]
        for shape in slide.shapes:
            if 'Title' in shape.name:
                title = shape
        
        if title:
            title.text = self.human.name + " Human Sequence Annotated"

    
    def populate_str_align_slide(self, align_imgs: list):
        """
        Populates the third slide of protein passport ppt template.

        Args:
            align_imgs (list): Structure align images.
            seq_img (Img): Aligned sequence image.
        """
        slide = self.slides[2]
        
        pictures = []
        captions = []
        for i in align_imgs:
            pictures.append(i.path)
            captions.append(i.caption)
            
        placeholders = []
        textboxes = []
        table = None
        title = None

        # Collect all shapes
        for shape in slide.shapes:
            if 'Picture' in shape.name:
                placeholders.append(shape)
            if 'Table' in shape.name:
                table = shape.table
            if 'TextBox' in shape.name:
                textboxes.append(shape)
            if 'Title' in shape.name:
                title = shape
        
        if title:
            if Annotation.ECD in self.human.annotations:
                title.text = 'ECD Alignment'
            else:
                title.text = 'Mature Alignment'
        
        # Insert pictures into placeholders (skip first placeholder if it's for something else)
        zipped = zip(placeholders[1:], pictures)
        for z in zipped:
            z[0].insert_picture(z[1])

        # Handle textboxes dynamically based on number of captions
        num_captions = len(captions)
        num_textboxes = len(textboxes)
        
        # If we have more textboxes than captions, delete the excess ones
        if num_textboxes > num_captions:
            textboxes_to_delete = textboxes[num_captions:]
            for textbox in textboxes_to_delete:
                # Delete the shape by removing its XML element
                textbox.element.getparent().remove(textbox.element)
            # Keep only the textboxes we need
            textboxes = textboxes[:num_captions]
        
        # If we have more captions than textboxes, create new ones
        elif num_captions > num_textboxes:
            # Calculate spacing and position based on existing textboxes
            if textboxes:
                last_textbox = textboxes[-1]
                # Try to determine spacing from existing textboxes
                if len(textboxes) > 1:
                    # Calculate spacing: find the minimum positive difference between consecutive textboxes
                    # This handles cases where textboxes might be in columns
                    top_positions = sorted(set([tb.top for tb in textboxes]))
                    if len(top_positions) > 1:
                        # Calculate differences between consecutive positions
                        diffs = [top_positions[i+1] - top_positions[i] for i in range(len(top_positions)-1)]
                        spacing = min([d for d in diffs if d > 0]) if diffs else last_textbox.height + Pt(10)
                    else:
                        # All textboxes at same vertical position (columns), use height + margin
                        spacing = last_textbox.height + Pt(10)
                else:
                    spacing = last_textbox.height + Pt(10)
                
                left = last_textbox.left
                top = last_textbox.top + spacing
                width = last_textbox.width
                height = last_textbox.height
            else:
                # Default position/size if no textboxes exist
                left = Inches(0.5)
                top = Inches(4)
                width = Inches(9)
                height = Inches(0.5)
                spacing = height + Pt(10)
            
            # Create new textboxes for remaining captions
            for i in range(num_textboxes, num_captions):
                new_top = top + (i - num_textboxes) * spacing
                new_textbox = slide.shapes.add_textbox(left, new_top, width, height)
                # Copy formatting from last textbox if available
                if textboxes:
                    last_textbox = textboxes[-1]
                    new_textbox.text_frame.word_wrap = last_textbox.text_frame.word_wrap
                    new_textbox.text_frame.margin_bottom = last_textbox.text_frame.margin_bottom
                    new_textbox.text_frame.margin_left = last_textbox.text_frame.margin_left
                    new_textbox.text_frame.margin_right = last_textbox.text_frame.margin_right
                    new_textbox.text_frame.margin_top = last_textbox.text_frame.margin_top
                textboxes.append(new_textbox)
        
        # Fill all textboxes with captions
        for textbox, caption in zip(textboxes, captions):
            textbox.text = caption
            for paragraph in textbox.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)

        if table:
            cell = table.cell(1, 1)
            cell.text = self.human.id
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

            for i, ortholog in enumerate(self.orthologs, start=2):
                species_cell = table.cell(i, 0)
                species_cell.text = ortholog.organism.name.capitalize()
                species_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

                id_cell = table.cell(i, 1)
                id_cell.text = ortholog.id
                id_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

                similarity_cell = table.cell(i, 2)
                similarity_cell.text = "%"
                similarity_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

        self.powerpoint.save(self.output_path)
    
    def populate_string_db_slide(self, network_img: str, pred_partners_img=None):
        """
        Populates the fourth slide of protein passport ppt template.

        Args:
            network_img (str): path to STRING network image.
            pred_partners_img (str): path to STRING predicted partners.
        """
        slide = self.slides[3]
        
        slide.shapes.add_picture(network_img, left=1, top=1)

        self.powerpoint.save(self.output_path)
    
    